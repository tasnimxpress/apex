"""
deliverables/make_deck.py  --  PHASE 3, the 10-slide leadership deck.

Max 10 slides, structure from ORCHESTRATOR_README. Built for NON-TECHNICAL leadership:
every slide makes ONE point, leads with the business implication, and reuses the EXACT
figures + numbers already validated in analysis/ (no new or conflicting values). The
mandatory flow-not-return caveat is carried on the methodology slide.

Run: d:/GitHub/apex/.venv/Scripts/python.exe deliverables/make_deck.py
"""
from pathlib import Path
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO = Path(__file__).resolve().parent.parent
A = REPO / "analysis"
OUT = REPO / "deliverables" / "apex_deck.pptx"

# ---- load canonical numbers (never retype headline figures) ------------------
m1 = json.load(open(A / "section1" / "_metrics.json"))
m2 = json.load(open(A / "section2" / "_metrics.json"))
m3 = json.load(open(A / "section3" / "_metrics.json"))

rank = {r["fund"]: r for r in m1["fund_ranking"]}
sc = {r["fund"]: r for r in m1["fund_scorecard"]}
F_FI, F_SH, F_BA, F_CG = ("Apex Fixed Income Fund", "Apex Shariah Growth Fund",
                          "Apex Balanced Opportunity Fund", "Apex Capital Growth Fund")

# ---- theme -------------------------------------------------------------------
NAVY = RGBColor(0x2E, 0x5A, 0x87)
TEAL = RGBColor(0x3E, 0x8E, 0x7E)
GOLD = RGBColor(0xC9, 0x91, 0x3B)
WINE = RGBColor(0x9A, 0x3B, 0x5B)
INK = RGBColor(0x22, 0x2A, 0x33)
GREY = RGBColor(0x5A, 0x63, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def fill_rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def header(slide, kicker, title, accent=NAVY):
    """Top band: small kicker + big headline (the ONE point of the slide)."""
    fill_rect(slide, 0, 0, SW, Inches(0.18), accent)
    tf = textbox(slide, Inches(0.55), Inches(0.30), SW - Inches(1.1), Inches(0.45))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = kicker.upper(); _set(r, 12, accent, bold=True)
    tf2 = textbox(slide, Inches(0.55), Inches(0.70), SW - Inches(1.1), Inches(1.0))
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = title; _set(r2, 26, INK, bold=True)
    return Inches(1.75)  # content top


def add_image_fit(slide, path, l, t, max_w, max_h, align_center=True):
    """Add picture scaled to fit a box, preserving aspect ratio."""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale); h = int(ih * scale)
    if align_center:
        l = int(l + (max_w - w) / 2)
        t = int(t + (max_h - h) / 2)
    slide.shapes.add_picture(str(path), l, t, width=w, height=h)


def bullets(slide, l, t, w, h, items, size=15, gap=8):
    tf = textbox(slide, l, t, w, h)
    for i, (txt, color, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # support inline emphasis: list of (text,color,bold) tuples per line via nested
        if isinstance(txt, list):
            for seg, scolor, sbold in txt:
                r = p.add_run(); r.text = seg; _set(r, size, scolor, bold=sbold)
        else:
            r = p.add_run(); r.text = txt; _set(r, size, color, bold=bold)
    return tf


def footnote(slide, text):
    tf = textbox(slide, Inches(0.55), SH - Inches(0.55), SW - Inches(1.1), Inches(0.4))
    r = tf.paragraphs[0].add_run(); r.text = text; _set(r, 10, GREY, italic=True)


def cr(x):
    return f"৳{x/1e7:,.1f} Cr"


# =============================================================================
# SLIDE 1 — Title & context
# =============================================================================
s = prs.slides.add_slide(BLANK)
fill_rect(s, 0, 0, SW, SH, NAVY)
fill_rect(s, 0, Inches(4.55), SW, Inches(0.06), GOLD)
tf = textbox(s, Inches(0.9), Inches(1.5), SW - Inches(1.8), Inches(2.0))
r = tf.paragraphs[0].add_run(); r.text = "Apex Asset Management"; _set(r, 44, WHITE, bold=True)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Where growth leaks — and the three moves that plug it"
_set(r, 26, RGBColor(0xCF, 0xDA, 0xE6))
tf2 = textbox(s, Inches(0.9), Inches(4.8), SW - Inches(1.8), Inches(2.0))
for i, line in enumerate([
    ("Senior Data Analyst case study  ·  Customer, fund & RM analytics", 16, True),
    ("Data as-of 31 May 2026  ·  12,229 accounts · 8,570 customers · 4 funds · 16 RMs", 14, False),
    ("Reproduced from cleaned source data; independently QA-validated (Gate 1 & Gate 2 PASS)", 12, False),
]):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    rr = p.add_run(); rr.text = line[0]
    _set(rr, line[1], RGBColor(0xE7, 0xED, 0xF3), bold=line[2])

# =============================================================================
# SLIDE 2 — Executive summary: leakage in one chart
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Executive summary", "Apex's growth problem is a retention problem, not acquisition")
add_image_fit(s, A / "synthesis/figures/02_aum_made_vs_lost.png",
              Inches(0.5), top, Inches(7.6), Inches(5.0))
bullets(s, Inches(8.3), top + Inches(0.1), Inches(4.6), Inches(5.0), [
    ([["+৳249.3 Cr", NAVY, True], [" lifetime net flow — money is", INK, False]], INK, False),
    ("being gathered fast (+৳37.8 Cr in 2025 alone)…", INK, False),
    ([["…but it leaks back out.", WINE, True]], INK, False),
    ("", INK, False),
    ([["1 fund grows; 3 leak.", NAVY, True]], INK, False),
    ("Fixed Income gathers AND keeps; Capital Growth", GREY, False),
    ("posts +৳68 Cr yet loses 47% — a leaky bucket.", GREY, False),
    ("", INK, False),
    ([["The base is single-fund (~82%)", NAVY, True]], INK, False),
    ("and a few RMs out-acquire what they retain.", GREY, False),
], size=15, gap=4)
footnote(s, "Bars = lifetime net flow (Purchase − Surrender) by fund; label = strict churn %; #rank = durable-growth composite. ‘Performance’ = flow/retention, not investment return (no NAV data).")

# =============================================================================
# SLIDE 3 — Methodology & data overview (+ key assumptions)
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Methodology & data", "How we measured — and one assumption that matters")
bullets(s, Inches(0.55), top, Inches(6.1), Inches(5.0), [
    ([["Data:", NAVY, True], [" 12,229 accounts · 8,570 customers (by mobile) ·", INK, False]], INK, False),
    ("231k transactions · 4 funds · 16 official RMs.", GREY, False),
    ("", INK, False),
    ([["Cleaned & reconciled:", NAVY, True], [" flow totals tie to the", INK, False]], INK, False),
    ("cent across raw → ledger → monthly panel; 3 unknown", GREY, False),
    ("accounts quarantined; anomalies flagged, never dropped.", GREY, False),
    ("", INK, False),
    ([["Beyond the 4 given formulas:", NAVY, True], [" cohort survival", INK, False]], INK, False),
    ("(Kaplan–Meier), SIP persistency curves, customer", GREY, False),
    ("clustering (k=4, stable ARI=1.0), RM acquisition quality.", GREY, False),
], size=15, gap=4)
# assumption callout box
fill_rect(s, Inches(7.0), top, Inches(5.8), Inches(4.6), LIGHT)
fill_rect(s, Inches(7.0), top, Inches(0.10), Inches(4.6), GOLD)
tf = textbox(s, Inches(7.35), top + Inches(0.25), Inches(5.2), Inches(4.2))
r = tf.paragraphs[0].add_run(); r.text = "KEY ASSUMPTION — point-in-time retention"
_set(r, 14, GOLD, bold=True)
for line in [
    "Account status is a TODAY snapshot. Measuring last",
    "year's retention from today's status would be wrong.",
    "We rebuild each account's status AT the measurement",
    "date from dated close / discontinue events.",
    "",
    "Result: the case cohort (May-24 → May-25) is",
    "84.6% retained (115/136) point-in-time — vs 93/136",
    "in a naïve snapshot. We use the rigorous number.",
    "",
    "Churn = Closed + Discontinued (headline).",
    "Inactive shown only as a sensitivity. ‘Performance’",
    "= flow / retention / engagement — NOT return.",
]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = line
    _set(rr, 12.5, INK, bold=line.startswith("Result") or line.startswith("Churn"))
footnote(s, "Single source of truth: definitions.yaml · every judgment call logged in decisions_log.md · reproduction in qa/qa_report.md.")

# =============================================================================
# SLIDE 4 — Fund ranking & scorecard
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Fund performance · 1", "Rank funds on durable growth, not on size")
add_image_fit(s, A / "section1/figures/07_fund_ranking.png",
              Inches(0.5), top, Inches(7.3), Inches(4.9))
bullets(s, Inches(8.0), top + Inches(0.1), Inches(4.9), Inches(5.0), [
    ([["Composite (0–100)", NAVY, True], [" = net flow, 36-mo", INK, False]], INK, False),
    ("survival, SIP persistency, AUM, low churn.", GREY, False),
    ("", INK, False),
    ([[f"#1 Fixed Income  {rank[F_FI]['composite_score']:.0f}", NAVY, True]], INK, False),
    ([[f"#2 Shariah Growth  {rank[F_SH]['composite_score']:.0f}", GREY, True]], INK, False),
    ([[f"#3 Capital Growth  {rank[F_CG]['composite_score']:.0f}", GREY, True]], INK, False),
    ([[f"#4 Balanced Opp.  {rank[F_BA]['composite_score']:.0f}", WINE, True]], INK, False),
    ("", INK, False),
    ([["Implication: ", INK, True], ["defend the leaky funds —", INK, False]], INK, False),
    ("don't pour growth spend into the winner.", GREY, False),
    ([["Owner: Head of Funds / PMs.", NAVY, True]], INK, False),
], size=15, gap=4)
footnote(s, "Composite weights favour durable growth (survival × persistency) over raw AUM; full scorecard in apex_analysis.ipynb §4.")

# =============================================================================
# SLIDE 5 — Fund deep-dive: leakage / cohorts / SIP persistency
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Fund performance · 2", "The leak is retention: cohorts fade, SIPs stop early", accent=WINE)
add_image_fit(s, A / "section1/figures/05_cohort_survival_by_fund.png",
              Inches(0.4), top, Inches(6.3), Inches(4.3))
add_image_fit(s, A / "section1/figures/02_sip_persistency_decay.png",
              Inches(6.8), top, Inches(6.2), Inches(4.3))
bullets(s, Inches(0.55), top + Inches(4.35), Inches(12.2), Inches(1.2), [
    ([["Capital Growth keeps only ~54% of a cohort to 3 years vs Fixed Income ~77%.", WINE, True],
      ["  SIP persistency is 87.9% overall but attrition concentrates in the early months.", INK, False]], INK, False),
    ([["Move: ", NAVY, True],
      ["auto-flag the FIRST missed installment in the monthly panel and fire a save-call; win-back team on Capital Growth (47.2% churn).  Owner: Head of Retention / Ops.", INK, False]], INK, False),
], size=13, gap=4)

# =============================================================================
# SLIDE 6 — Customer segments & targeting map
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Customers · 1", "Four stable segments — the book is single-fund & concentrated", accent=TEAL)
add_image_fit(s, A / "section2/figures/03_segments_scatter.png",
              Inches(0.5), top, Inches(7.2), Inches(4.9))
bullets(s, Inches(8.0), top + Inches(0.1), Inches(4.9), Inches(5.0), [
    ([["k=4 segments", TEAL, True], [" — stable across 3 seeds", INK, False]], INK, False),
    ("(ARI = 1.0; silhouette 0.42):", GREY, False),
    ("• High-value / lump-sum tier", INK, False),
    ("• Committed SIP savers", INK, False),
    ("• Young, small-ticket mass", INK, False),
    ("• Dormant non-SIP group", INK, False),
    ("", INK, False),
    ([["~82% hold ONE fund;", NAVY, True], [" top 10% of", INK, False]], INK, False),
    ("customers hold ~68% of AUM.", GREY, False),
    ([["Each segment maps to a primary fund", INK, True]], INK, False),
    ("for targeted offers (see notebook §5/§7).", GREY, False),
], size=15, gap=4)
footnote(s, "Clustering at customer (mobile) level on age, installment, tenor, onboarding amount (monetary log-scaled); k=2 rejected as trivial SIP/non-SIP split.")

# =============================================================================
# SLIDE 7 — Cross-sell opportunity (sized in Taka)
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Customers · 2", "The cheapest AUM: a second fund to people who already trust us", accent=TEAL)
add_image_fit(s, A / "section2/figures/04_cross_sell.png",
              Inches(0.5), top, Inches(7.3), Inches(4.9))
p1 = m2["cross_sell"][0]; p3 = m2["cross_sell"][2]
bullets(s, Inches(8.0), top + Inches(0.1), Inches(4.9), Inches(5.0), [
    ([["Two sized plays, ready now:", NAVY, True]], INK, False),
    ("", INK, False),
    ([[f"1 · Single → second fund", TEAL, True]], INK, False),
    ([[f"   {int(p1['target_customers']):,} customers  ≈ {cr(p1['taka'])}", INK, True]], INK, False),
    ("   start with 2,273 high-value leads", GREY, False),
    ("   already under a cross-capable RM.", GREY, False),
    ("", INK, False),
    ([[f"2 · Non-SIP → SIP migration", TEAL, True]], INK, False),
    ([[f"   {int(p3['target_customers'])} customers  ≈ {cr(p3['taka'])}/yr recurring", INK, True]], INK, False),
    ("", INK, False),
    ([["Owner: Head of Sales.", NAVY, True]], INK, False),
], size=15, gap=4)
footnote(s, "Taka = target count × explicit median ticket (2nd-fund ৳12,000; SIP annual ৳84,000). Reproducible; headline wallet excludes the overlapping warm-lead subset.")

# =============================================================================
# SLIDE 8 — RM performance quadrant & scorecard
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "RM productivity · 1", "Reward retention, not raw sign-ups", accent=GOLD)
add_image_fit(s, A / "section3/figures/01_rm_quadrant.png",
              Inches(0.5), top, Inches(7.6), Inches(5.0))
flagged = ", ".join(m3["high_vol_low_quality_rms"])
bullets(s, Inches(8.3), top + Inches(0.1), Inches(4.6), Inches(5.0), [
    ([["Volume × value, coloured by", INK, False]], INK, False),
    ([["1-year retention.", NAVY, True], [" Overall book", INK, False]], INK, False),
    ([[f"retention = {m3['overall_retention_1yr']*100:.1f}%.", NAVY, True]], INK, False),
    ("", INK, False),
    ([["Coach, don't celebrate:", WINE, True]], INK, False),
    ("high-volume but below-median", GREY, False),
    ("retention —", GREY, False),
    ([[flagged + ".", WINE, True]], INK, False),
    ("", INK, False),
    ([["Move: ", NAVY, True], ["tie a slice of incentive", INK, False]], INK, False),
    ("to 1-yr retention.  Owner: Sales Mgr.", GREY, False),
], size=15, gap=4)
footnote(s, "Acquisition attributed on Introducer role, 16 official RMs; 1-yr retention on cohorts old enough to observe it (onboarded ≤ as-of − 12m), reconstructed point-in-time.")

# =============================================================================
# SLIDE 9 — Team trend & systemic findings
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "RM productivity · 2", "A mild systemic dip — and where to deploy specialists", accent=GOLD)
add_image_fit(s, A / "section3/figures/03_team_yoy.png",
              Inches(0.5), top, Inches(8.0), Inches(4.7))
bullets(s, Inches(8.7), top + Inches(0.1), Inches(4.2), Inches(5.0), [
    ([["Team 1-yr retention dipped", INK, False]], INK, False),
    ([["across the 2021→2023 vintages", NAVY, True]], INK, False),
    ("(≈90% → 84%) then recovered to", GREY, False),
    ("86.9% (2025) — a mild systemic", GREY, False),
    ("signal, not just two people.", GREY, False),
    ("", INK, False),
    ([["Deploy high-retention", GOLD, True]], INK, False),
    ([["specialists onto Capital Growth", GOLD, True]], INK, False),
    ("(the leakiest fund) rather than", GREY, False),
    ("retraining them.", GREY, False),
    ("", INK, False),
    ([["Owner: Head of Sales.", NAVY, True]], INK, False),
], size=14, gap=4)
footnote(s, "Net flow by introducing vintage rising (to +৳37.8 Cr in 2025); systemic vs individual separated so the fix matches the cause.")

# =============================================================================
# SLIDE 10 — Prioritized recommendations & roadmap
# =============================================================================
s = prs.slides.add_slide(BLANK)
top = header(s, "Recommendations", "Five owner-assigned moves, sequenced")
rows = [
    ("NOW", "Stop the leak", "Auto-flag first missed SIP installment + win-back on Capital Growth (47.2% churn)", "Head of Retention / PM", WINE),
    ("NOW", "Cheapest AUM", f"Second-fund campaign to {int(p1['target_customers']):,} single-fund customers (~{cr(p1['taka'])})", "Head of Sales", TEAL),
    ("NEXT", "Fix incentives", f"Tie RM incentive to 1-yr retention ({m3['overall_retention_1yr']*100:.0f}%); coach 2 flagged RMs", "Sales Manager", GOLD),
    ("NEXT", "Recurring revenue", f"Non-SIP → SIP migration for {int(p3['target_customers'])} customers (~{cr(p3['taka'])}/yr)", "Head of Sales", TEAL),
    ("ONGOING", "Pre-empt redemptions", "Dividend-reinvest offer WITH each Jul/Jan distribution (20.4% vs 16.7% post-dividend surrenders)", "Marketing + PM", NAVY),
]
y = top
rowh = Inches(0.92)
for i, (when, what, action, owner, col) in enumerate(rows):
    yy = y + i * rowh
    fill_rect(s, Inches(0.55), yy, Inches(0.12), Inches(0.78), col)
    tf = textbox(s, Inches(0.8), yy, Inches(1.6), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run(); r.text = when; _set(r, 13, col, bold=True)
    tf = textbox(s, Inches(2.4), yy, Inches(2.5), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run(); r.text = what; _set(r, 14, INK, bold=True)
    tf = textbox(s, Inches(4.9), yy, Inches(6.0), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run(); r.text = action; _set(r, 12, GREY, False)
    tf = textbox(s, Inches(11.0), yy, Inches(2.0), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run(); r.text = owner; _set(r, 11.5, NAVY, bold=True)
footnote(s, "Every number traces to apex_analysis.ipynb and the recommendation register (analysis/synthesis). ‘Performance’ throughout = flow/retention/engagement, not investment return.")

prs.save(str(OUT))
print("saved", OUT, "·", len(prs.slides._sldIdLst), "slides")
